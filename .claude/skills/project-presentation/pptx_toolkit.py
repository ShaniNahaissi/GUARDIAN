#!/usr/bin/env python3
"""Reusable python-pptx helper toolkit for project/architecture decks.

Copy this file into the scratchpad and import from it (or paste the parts you
need directly into your build script) rather than re-deriving these helpers —
they already work around python-pptx's sharpest edge (see CRITICAL below).

Usage sketch:

    from pptx_toolkit import *
    prs, add_slide, set_bg = new_deck()
    s = content_slide(prs, add_slide, "Architecture", "System overview", page_no=3)
    flow_row(s, [("Step 1", "detail"), ("Step 2", "detail")], y0=Inches(2.5),
              box_h=Inches(1.5), box_w=Inches(2.5), gap=Inches(0.3))
    prs.save("out.pptx")
    validate_pptx("out.pptx")   # <- always run this before calling it done

CRITICAL — read before writing any shape/connector coordinate math:
Every x/y/cx/cy fed to a python-pptx shape or connector MUST be an int. `Emu`
is an int subclass, but Python 3's `/` operator always returns a float even
on int subclasses (`Emu(100) / 2 == 50.0`, a float). If that float reaches
shape XML (e.g. `x="2909163.5"`), PowerPoint treats the file as CORRUPT and
shows a repair prompt on open — this is not cosmetic, the file is genuinely
invalid OOXML. Rules:
  - Never write `something / 2` for a coordinate. Use `half()` below or `// 2`.
  - Every helper in this file wraps its x/y/w/h args in `Emu(int(...))`
    defensively — use these helpers rather than calling
    `slide.shapes.add_shape(...)` / `add_connector(...)` directly.
  - After saving, ALWAYS run `validate_pptx()` (bottom of this file) — it
    reloads the file and scans every slide's raw XML for stray float
    coordinates. Zero matches is the only acceptable result.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
import re
import zipfile

# --- Default brand palette — swap these for the target project/company colors ---
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_LIGHT = RGBColor(0x14, 0x33, 0x5C)
TEAL = RGBColor(0x1A, 0xAE, 0x9F)
AMBER = RGBColor(0xE8, 0xA2, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
MID_GRAY = RGBColor(0x6B, 0x76, 0x85)
DARK_TEXT = RGBColor(0x1B, 0x22, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER_BG = RGBColor(0xFB, 0xF2, 0xE3)
AMBER_TEXT = RGBColor(0x8A, 0x5A, 0x12)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9


def half(emu):
    """Integer half of an Emu length. Use this instead of `emu / 2`."""
    return Emu(int(emu) // 2)


def new_deck():
    """Returns a fresh 16:9 Presentation with a blank-layout add_slide() bound."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    return prs, add_slide


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, rounded=False, line_color=None, line_w=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=DARK_TEXT, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=1.0):
    """text may contain \\n for multiple paragraphs."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
    return tb


def add_box_label(slide, x, y, w, h, text, size=13, color=WHITE, bold=True,
                   align=PP_ALIGN.CENTER, line_spacing=1.05):
    """Centered label meant to sit INSIDE a shape you already drew with add_rect."""
    return add_text(slide, x, y, w, h, text, size=size, color=color, bold=bold,
                     align=align, anchor=MSO_ANCHOR.MIDDLE, line_spacing=line_spacing)


def add_arrow(slide, x1, y1, x2, y2, color=MID_GRAY, weight=2.25):
    """Straight connector with a triangle arrowhead at (x2, y2)."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)),
                                       Emu(int(x2)), Emu(int(y2)))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return conn


def add_footer(slide, page_no, label="Internal Briefing"):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), NAVY)
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.32), Inches(8), Inches(0.32), label,
              size=9, color=RGBColor(0xB9, 0xC4, 0xD3), anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.32), Inches(0.9), Inches(0.32),
              str(page_no), size=9, color=RGBColor(0xB9, 0xC4, 0xD3), align=PP_ALIGN.RIGHT,
              anchor=MSO_ANCHOR.MIDDLE)


def content_slide(add_slide, kicker, title, page_no, footer_label="Internal Briefing"):
    """Standard header-bar content slide. Returns the slide — add your own body content."""
    s = add_slide()
    set_bg(s, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_rect(s, 0, Inches(1.15), SLIDE_W, Inches(0.06), TEAL)
    add_text(s, Inches(0.55), Inches(0.14), Inches(10), Inches(0.35), kicker.upper(),
              size=12, color=TEAL, bold=True)
    add_text(s, Inches(0.55), Inches(0.42), Inches(12), Inches(0.65), title,
              size=26, color=WHITE, bold=True)
    add_footer(s, page_no, footer_label)
    return s


def section_slide(add_slide, number, title, subtitle, page_no, footer_label="Internal Briefing"):
    """Full-bleed navy divider slide for a new section."""
    s = add_slide()
    set_bg(s, NAVY)
    add_rect(s, Inches(0.9), Inches(3.35), Inches(1.6), Inches(0.06), TEAL)
    add_text(s, Inches(0.9), Inches(2.5), Inches(3), Inches(1), number,
              size=70, color=NAVY_LIGHT, bold=True)
    add_text(s, Inches(0.9), Inches(3.5), Inches(11), Inches(1.1), title,
              size=40, color=WHITE, bold=True)
    add_text(s, Inches(0.9), Inches(4.35), Inches(11), Inches(0.8), subtitle,
              size=16, color=RGBColor(0xB9, 0xC4, 0xD3))
    add_footer(s, page_no, footer_label)
    return s


def title_slide(add_slide, kicker, title, subtitle, tagline=None):
    """Full-bleed navy title/cover slide. `title` may contain \\n."""
    s = add_slide()
    set_bg(s, NAVY)
    add_rect(s, 0, Inches(4.9), SLIDE_W, Inches(0.05), TEAL)
    add_text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.5), kicker,
              size=20, color=TEAL, bold=True)
    add_text(s, Inches(0.85), Inches(3.05), Inches(11.5), Inches(1.5), title,
              size=40, color=WHITE, bold=True, line_spacing=1.05)
    add_text(s, Inches(0.9), Inches(5.1), Inches(11), Inches(0.5), subtitle,
              size=17, color=RGBColor(0xB9, 0xC4, 0xD3))
    if tagline:
        add_text(s, Inches(0.9), Inches(6.75), Inches(11), Inches(0.4), tagline,
                  size=12, color=MID_GRAY)
    return s


def flow_row(slide, steps, y0, box_h, box_w, gap, default_color=NAVY_LIGHT,
             highlight_index=None, highlight_color=TEAL, label_size=14, desc_size=None,
             text_color=WHITE):
    """Horizontal row of connected flow boxes, auto-centered on the slide width.

    `steps`: list of (title, desc) or (title, desc, color) — title/desc may use \\n.
    desc may be None for a plain label-only box (no secondary text).
    Returns start_x (the x of the first box), useful for annotating below the row.
    """
    n = len(steps)
    total = box_w * n + gap * (n - 1)
    start_x = (SLIDE_W - total) // 2
    for i, step in enumerate(steps):
        title = step[0]
        desc = step[1] if len(step) > 1 else None
        color = step[2] if len(step) > 2 else (highlight_color if i == highlight_index else default_color)
        x = start_x + (box_w + gap) * i
        add_rect(slide, x, y0, box_w, box_h, color, rounded=True)
        if desc:
            add_box_label(slide, x + Inches(0.1), y0 + Inches(0.15), box_w - Inches(0.2), Inches(0.65),
                           title, size=label_size, color=text_color, bold=True)
            add_text(slide, x + Inches(0.15), y0 + Inches(0.85), box_w - Inches(0.3), box_h - Inches(0.95),
                      desc, size=desc_size or (label_size - 4), color=RGBColor(0xDB, 0xE2, 0xEB),
                      align=PP_ALIGN.CENTER, line_spacing=1.1)
        else:
            add_box_label(slide, x + Inches(0.05), y0, box_w - Inches(0.1), box_h, title,
                           size=label_size, color=text_color, bold=True)
        if i < n - 1:
            add_arrow(slide, x + box_w + Inches(0.03), y0 + half(box_h),
                      x + box_w + gap - Inches(0.03), y0 + half(box_h), MID_GRAY, 2.25)
    return start_x


def metric_placeholder_row(slide, metrics, y0=Inches(2.5), box_h=Inches(2.7)):
    """Row of "TBD" metric cards for numbers not yet available (model/app performance).

    `metrics`: list of (label, description) tuples.
    """
    n = len(metrics)
    box_w = Inches(2.7); gap = Inches(0.35)
    total = box_w * n + gap * (n - 1)
    start_x = (SLIDE_W - total) // 2
    for i, (label, desc) in enumerate(metrics):
        x = start_x + (box_w + gap) * i
        add_rect(slide, x, y0, box_w, box_h, LIGHT_GRAY, rounded=True)
        add_rect(slide, x, y0, box_w, Inches(0.12), TEAL)
        add_text(slide, x + Inches(0.15), y0 + Inches(0.35), box_w - Inches(0.3), Inches(1.1), "TBD",
                  size=40, color=NAVY_LIGHT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y0 + Inches(1.55), box_w - Inches(0.3), Inches(0.4), label,
                  size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y0 + Inches(2.0), box_w - Inches(0.3), Inches(0.65), desc,
                  size=11, color=MID_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.1)


def callout_banner(slide, x, y, w, h, text, bg=AMBER_BG, bar=AMBER, text_color=AMBER_TEXT, size=13.5):
    """Amber 'heads up / placeholder / status' banner — used for TBD metrics or known gaps."""
    add_rect(slide, x, y, w, h, bg, rounded=True)
    add_rect(slide, x, y, Inches(0.1), h, bar)
    add_text(slide, x + Inches(0.35), y, w - Inches(0.5), h, text, size=size, color=text_color,
              bold=True, anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, header, rows, col_widths=None,
              header_bg=NAVY, header_fg=WHITE, zebra=LIGHT_GRAY):
    """Simple styled table. `header`: list[str]. `rows`: list[list[str]]."""
    n_rows, n_cols = len(rows) + 1, len(header)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Emu(int(cw))
    for c, htext in enumerate(header):
        cell = tbl.cell(0, c)
        cell.text = htext
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.color.rgb = header_fg; p.runs[0].font.bold = True; p.runs[0].font.size = Pt(13)
    for r, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(r, ci)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = zebra if r % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = DARK_TEXT
            p.runs[0].font.bold = (ci == 0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
    return tbl


def add_bullets(slide, x, y, w, h, items, size=15, color=DARK_TEXT, bullet_color=TEAL,
                 space_after=10, line_spacing=1.08):
    """items: list of (text, level) or (text, level, bold_lead_substring) or plain str (level 0)."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item[0], item[1]
            lead = item[2] if len(item) > 2 else None
        else:
            text, level, lead = item, 0, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        msize = size if level == 0 else size - 1
        mcolor = color if level == 0 else MID_GRAY
        r0 = p.add_run()
        r0.text = "▸  " if level == 0 else "–  "
        r0.font.size = Pt(msize); r0.font.bold = True
        r0.font.color.rgb = bullet_color if level == 0 else MID_GRAY
        if lead:
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(msize); r1.font.bold = True; r1.font.color.rgb = NAVY
            rest = text[len(lead):]
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.size = Pt(msize); r2.font.color.rgb = mcolor
        else:
            r1 = p.add_run(); r1.text = text
            r1.font.size = Pt(msize); r1.font.color.rgb = mcolor
    return tb


def validate_pptx(path):
    """Reload the saved file and scan for stray float EMU coordinates.

    ALWAYS call this after prs.save(path) — see the CRITICAL note at the top
    of this file. Raises AssertionError if anything looks corrupt; otherwise
    prints a one-line OK summary.
    """
    Presentation(path)  # raises if python-pptx itself can't parse it
    z = zipfile.ZipFile(path)
    bad = []
    for name in z.namelist():
        if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
            data = z.read(name).decode("utf-8")
            bad += [(name, m.group(0)) for m in re.finditer(r'(?:x|y|cx|cy)="(-?\d+\.\d+)"', data)]
    assert not bad, f"Found {len(bad)} float EMU coordinates (will corrupt in PowerPoint): {bad[:10]}"
    n_slides = len(Presentation(path).slides._sldIdLst)
    print(f"OK: {path} — {n_slides} slides, reloads cleanly, no float-EMU coordinates.")
